# add_slides.py - Appends missing slides to StegAI_SDP_Presentation.pptx
# Run: venv/Scripts/python add_slides.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (match existing) ────────────────────────────────────────────
BG      = RGBColor(0x0F,0x17,0x2A)
ACCENT  = RGBColor(0x38,0xBD,0xF8)
ACCENT2 = RGBColor(0x34,0xD3,0x99)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
LIGHT   = RGBColor(0xCB,0xD5,0xE1)
YELLOW  = RGBColor(0xFB,0xBF,0x24)
CARD_BG = RGBColor(0x1E,0x29,0x3B)
RED     = RGBColor(0xF4,0x72,0x6D)
PURPLE  = RGBColor(0xA7,0x8B,0xFA)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation("StegAI_SDP_Presentation.pptx")
blank = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────────────
def add_bg(sl):
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = BG

def rect(sl,l,t,w,h,col):
    s=sl.shapes.add_shape(1,l,t,w,h); s.fill.solid()
    s.fill.fore_color.rgb=col; s.line.fill.background(); return s

def txt(sl,text,l,t,w,h,size=15,bold=False,col=WHITE,align=PP_ALIGN.LEFT,italic=False):
    bx=sl.shapes.add_textbox(l,t,w,h); bx.word_wrap=True
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=col; r.font.italic=italic
    return bx

def bullets(sl,items,l,t,w,h,size=14,col=LIGHT,marker="▸ "):
    bx=sl.shapes.add_textbox(l,t,w,h); bx.word_wrap=True
    tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(4)
        r=p.add_run(); r.text=marker+item
        r.font.size=Pt(size); r.font.color.rgb=col

def header(sl,title,sub=None):
    rect(sl,0,0,SLIDE_W,Inches(0.08),ACCENT)
    txt(sl,title,Inches(0.5),Inches(0.15),Inches(12),Inches(0.72),
        size=30,bold=True,col=WHITE)
    if sub:
        txt(sl,sub,Inches(0.5),Inches(0.82),Inches(12),Inches(0.4),
            size=14,col=ACCENT,italic=True)
    rect(sl,0,Inches(7.3),SLIDE_W,Inches(0.2),CARD_BG)

def card(sl,l,t,w,h):
    rect(sl,l,t,w,h,CARD_BG)

# ════════════════════════════════════════════════════════════════════════
# SLIDE A – Existing vs Proposed System
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Existing vs Proposed System","What's Different About StegAI?")

card(sl,Inches(0.4),Inches(1.4),Inches(5.9),Inches(5.6))
txt(sl,"❌  Existing Approach (Traditional LSB)",
    Inches(0.6),Inches(1.55),Inches(5.5),Inches(0.45),
    size=17,bold=True,col=RED)
bullets(sl,[
    "Basic LSB — data embedded uniformly across all pixels",
    "Ignores image texture, complexity, and region suitability",
    "Causes predictable statistical patterns in pixel values",
    "Easily detected by: Histogram Analysis, Chi-square tests, SPAM/SRM models",
    "No encryption — data exposed even if stego image is found",
    "No payload control — risk of over-embedding",
],Inches(0.6),Inches(2.1),Inches(5.5),Inches(4.7),size=14)

card(sl,Inches(6.65),Inches(1.4),Inches(6.3),Inches(5.6))
txt(sl,"✅  Proposed System (StegAI)",
    Inches(6.85),Inches(1.55),Inches(5.9),Inches(0.45),
    size=17,bold=True,col=ACCENT2)
bullets(sl,[
    "Image divided into 8×8 blocks for analysis",
    "Each block evaluated using Entropy + Variance",
    "High-texture regions selected as safe embedding zones",
    "Smooth/flat regions are avoided entirely",
    "Adaptive LSB — embedding locations are controlled",
    "AES-256 encrypts message BEFORE hiding it",
    "Reduces statistical distortion and steganalysis detectability",
],Inches(6.85),Inches(2.1),Inches(5.9),Inches(4.7),size=14)

# ════════════════════════════════════════════════════════════════════════
# SLIDE B – Literature Survey (Part 1)
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Literature Survey","Related Work & Prior Art — Part 1 of 2")

papers1=[
    ("Pevný et al., 2010","SPAM","Subtractive Pixel Adjacency Matrix to detect stego artifacts via pixel differences",
     "Effective against LSB; low compute cost","Limited against modern adaptive steganography",ACCENT),
    ("Fridrich & Kodovský, 2012","SRM","Rich Models using high-dimensional residual-based features for steganalysis",
     "Very high detection accuracy; benchmark standard","Large feature space; not suitable for real-time",ACCENT2),
    ("Yang et al., 2023","IStego100K","Large-scale benchmark dataset for training & evaluating steganalysis models",
     "Diverse embedding methods; supports ML training","High storage required; not for lightweight systems",YELLOW),
]

for i,(authors,name,desc,pro,con,col) in enumerate(papers1):
    ty=Inches(1.5+i*1.9)
    card(sl,Inches(0.4),ty,Inches(12.5),Inches(1.75))
    rect(sl,Inches(0.4),ty,Inches(2.2),Inches(1.75),col)
    txt(sl,name,Inches(0.45),ty+Inches(0.15),Inches(2.1),Inches(0.5),
        size=16,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,authors,Inches(0.45),ty+Inches(0.65),Inches(2.1),Inches(0.5),
        size=11,col=BG,align=PP_ALIGN.CENTER,italic=True)
    txt(sl,desc,Inches(2.75),ty+Inches(0.1),Inches(4.4),Inches(1.5),size=13,col=LIGHT)
    txt(sl,"✅ "+pro,Inches(7.3),ty+Inches(0.1),Inches(2.7),Inches(0.7),size=12,col=ACCENT2)
    txt(sl,"❌ "+con,Inches(7.3),ty+Inches(0.85),Inches(2.7),Inches(0.7),size=12,col=RED)

# ════════════════════════════════════════════════════════════════════════
# SLIDE C – Literature Survey (Part 2)
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Literature Survey","Related Work & Prior Art — Part 2 of 2")

papers2=[
    ("Holub et al., 2014","UNIWARD","Universal distortion function for content-adaptive steganography across domains",
     "Reduced detectability; adaptable to multiple domains","Computationally intensive; complex distortion modeling",PURPLE),
    ("Apau et al., 2025","Survey 2025","Comprehensive survey of techniques resisting statistical steganalysis attacks",
     "Categorizes adaptive, statistical & content-aware methods","Survey only; no unified algorithm proposed",ACCENT),
    ("De La Croix, 2024","DL Survey","Survey on deep learning–based steganalysis — CNN and hybrid techniques",
     "Explains evolution from statistical to deep learning","Detection-focused; requires large datasets",ACCENT2),
]

for i,(authors,name,desc,pro,con,col) in enumerate(papers2):
    ty=Inches(1.5+i*1.9)
    card(sl,Inches(0.4),ty,Inches(12.5),Inches(1.75))
    rect(sl,Inches(0.4),ty,Inches(2.2),Inches(1.75),col)
    txt(sl,name,Inches(0.45),ty+Inches(0.15),Inches(2.1),Inches(0.5),
        size=16,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,authors,Inches(0.45),ty+Inches(0.65),Inches(2.1),Inches(0.5),
        size=11,col=BG,align=PP_ALIGN.CENTER,italic=True)
    txt(sl,desc,Inches(2.75),ty+Inches(0.1),Inches(4.4),Inches(1.5),size=13,col=LIGHT)
    txt(sl,"✅ "+pro,Inches(7.3),ty+Inches(0.1),Inches(2.7),Inches(0.7),size=12,col=ACCENT2)
    txt(sl,"❌ "+con,Inches(7.3),ty+Inches(0.85),Inches(2.7),Inches(0.7),size=12,col=RED)

# highlight gap
card(sl,Inches(10.2),Inches(1.5),Inches(2.7),Inches(5.7))
txt(sl,"Research\nGap",Inches(10.3),Inches(1.6),Inches(2.5),Inches(0.6),
    size=15,bold=True,col=YELLOW,align=PP_ALIGN.CENTER)
bullets(sl,[
    "No system combines ML + Crypto + Adaptive Steg",
    "No web-based real-time pipeline exists",
    "StegAI fills this gap",
],Inches(10.25),Inches(2.3),Inches(2.55),Inches(4.5),size=12,col=LIGHT,marker="→ ")

# ════════════════════════════════════════════════════════════════════════
# SLIDE D – Threat Model & Attack Analysis
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Threat Model & Attack Analysis","How StegAI Resists Detection")

txt(sl,"StegAI is designed to resist Low-Level and Mid-Level steganalysis attacks through AI-guided adaptive embedding and controlled payload capacity.",
    Inches(0.5),Inches(1.1),Inches(12.4),Inches(0.6),size=14,col=LIGHT,italic=True)

attacks=[
    ("Low-Level Attacks",ACCENT,[
        "LSB Statistical Analysis — detects unnatural LSB randomness",
        "Histogram Analysis — compares pixel intensity distributions",
        "Noise Variance Analysis — detects added noise from embedding",
    ]),
    ("Mid-Level Attacks",YELLOW,[
        "Chi-Square Attack — tests pixel pair probability distributions",
        "SPAM-Based Analysis — checks pixel adjacency & difference patterns",
        "Feature-Based ML Attacks — SVM/RF using entropy & co-occurrence",
        "Texture Inconsistency Detection — abnormal texture transitions",
    ]),
]

for i,(title,col,items) in enumerate(attacks):
    lx=Inches(0.4+i*6.55)
    card(sl,lx,Inches(1.85),Inches(6.2),Inches(5.2))
    rect(sl,lx,Inches(1.85),Inches(6.2),Inches(0.45),col)
    txt(sl,title,lx+Inches(0.1),Inches(1.85),Inches(6.0),Inches(0.45),
        size=16,bold=True,col=BG)
    bullets(sl,items,lx+Inches(0.15),Inches(2.4),Inches(5.9),Inches(4.3),size=14)

card(sl,Inches(0.4),Inches(7.1),Inches(12.5),Inches(0.45))
txt(sl,"StegAI Defence: Selects only high-entropy blocks  ▸  Limits payload per region  ▸  Encrypts before embedding  ▸  Histogram stays natural",
    Inches(0.6),Inches(7.1),Inches(12.2),Inches(0.45),size=13,col=ACCENT2,bold=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE E – Technology Stack
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Technology Stack","Tools, Libraries & Frameworks Used")

stack=[
    ("🖥  Frontend",    ACCENT,  ["HTML5 + CSS3","JavaScript (Vanilla)","Jinja2 Templates","Responsive UI — image preview + download"]),
    ("⚙️  Backend",     ACCENT2, ["Python 3.x","Flask (REST API)","Werkzeug (file security)","RESTful route handling"]),
    ("🔐  Encryption",  YELLOW,  ["cryptography library","ECC — SECP256R1 curve","HKDF-SHA256 key derivation","AES-256-CFB mode"]),
    ("🖼  Image Proc.", PURPLE,  ["OpenCV (cv2)","NumPy arrays","LSB steganography","PSNR / SSIM metrics"]),
    ("🤖  ML Module",   RED,     ["scikit-learn","Random Forest (50 trees)","Block feature extraction","std & variance features"]),
    ("📊  Evaluation",  ACCENT,  ["PSNR > 40 dB target","SSIM ≈ 1.0 target","Matplotlib histogram","Side-by-side comparison"]),
]

for i,(title,col,items) in enumerate(stack):
    lx=Inches(0.3 + (i%3)*4.35)
    ty=Inches(1.4 + (i//3)*2.9)
    card(sl,lx,ty,Inches(4.1),Inches(2.65))
    rect(sl,lx,ty,Inches(4.1),Inches(0.42),col)
    txt(sl,title,lx+Inches(0.1),ty,Inches(3.9),Inches(0.42),
        size=14,bold=True,col=BG)
    bullets(sl,items,lx+Inches(0.1),ty+Inches(0.48),Inches(3.9),Inches(2.1),
            size=13,marker="• ")

# ════════════════════════════════════════════════════════════════════════
# SLIDE F – Modules Description
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Modules Description","Five Core Components of StegAI")

modules=[
    ("Image Analysis\nModule",   ACCENT,  "ecc_module.py",
     ["Divides image into 8×8 pixel blocks","Computes Shannon entropy per block","Calculates variance for texture scoring","Identifies safe embedding regions"]),
    ("ECC Encryption\nModule",   ACCENT2, "ecc_module.py",
     ["Generates SECP256R1 key pair","ECDH + HKDF derives shared AES key","Encrypts message before embedding","Ensures security even if stego is found"]),
    ("Adaptive Embedding\nModule", YELLOW,"steg_module.py",
     ["Embeds only in ML-approved safe blocks","Controls bits per pixel dynamically","LSB embedding with 0xFE mask","Appends STEGO123 delimiter"]),
    ("Extraction &\nDecryption",  PURPLE, "steg_module.py",
     ["Reads LSBs sequentially from image","Assembles bytes until delimiter found","Strips AES key and IV from payload","Decrypts and returns plaintext"]),
    ("Evaluation\nModule",        RED,    "steg_module.py",
     ["Computes PSNR (target > 40 dB)","Computes SSIM (target ≈ 0.99+)","Generates pixel histogram comparison","Matplotlib — saves to static/hist.png"]),
]

positions=[(Inches(0.3),Inches(1.4)),(Inches(2.9),Inches(1.4)),(Inches(5.5),Inches(1.4)),
           (Inches(8.1),Inches(1.4)),(Inches(10.7),Inches(1.4))]
for (lx,ty),(name,col,file,items) in zip(positions,modules):
    card(sl,lx,ty,Inches(2.4),Inches(5.7))
    rect(sl,lx,ty,Inches(2.4),Inches(0.85),col)
    txt(sl,name,lx+Inches(0.1),ty+Inches(0.05),Inches(2.2),Inches(0.75),
        size=13,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,file,lx+Inches(0.1),ty+Inches(0.9),Inches(2.2),Inches(0.35),
        size=10,col=ACCENT,italic=True,align=PP_ALIGN.CENTER)
    bullets(sl,items,lx+Inches(0.1),ty+Inches(1.3),Inches(2.2),Inches(4.1),
            size=12,marker="→ ")

# ════════════════════════════════════════════════════════════════════════
# SLIDE G – Implementation Progress
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Implementation Progress","What's Done · What's In Progress")

done=[
    "Web-based interface — image upload, message input, preview & download",
    "ECC (SECP256R1) encryption and decryption — fully working",
    "AES-256-CFB encrypt / decrypt pipeline — verified",
    "Basic LSB embedding and extraction — tested on sample images",
    "AI feature extraction (entropy + variance) — completed",
    "Random Forest ML classifier — trained and integrated",
    "PSNR & SSIM metrics — computed and displayed in UI",
    "Pixel histogram comparison — auto-generated and shown",
    "End-to-end pipeline — tested successfully on multiple images",
]
inprog=[
    "Integration of full ECC ECDH two-party key exchange",
    "Frequency-domain features (DCT) for ML classifier",
    "Dynamic payload capacity control per block",
    "Chi-square steganalysis resistance benchmarking",
    "Performance testing on large image datasets",
]

card(sl,Inches(0.4),Inches(1.4),Inches(7.2),Inches(5.7))
txt(sl,"✅  Completed",Inches(0.6),Inches(1.55),Inches(6.8),Inches(0.45),
    size=18,bold=True,col=ACCENT2)
bullets(sl,done,Inches(0.6),Inches(2.1),Inches(6.8),Inches(4.8),size=14)

card(sl,Inches(7.85),Inches(1.4),Inches(5.1),Inches(5.7))
txt(sl,"🔄  In Progress",Inches(8.0),Inches(1.55),Inches(4.8),Inches(0.45),
    size=18,bold=True,col=YELLOW)
bullets(sl,inprog,Inches(8.0),Inches(2.1),Inches(4.8),Inches(4.8),size=14)

# ════════════════════════════════════════════════════════════════════════
# SLIDE H – Evaluation Strategy
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"Evaluation Strategy","Measuring Quality, Security & Performance")

evals=[
    ("📈  Image Quality",ACCENT,[
        "PSNR — Peak Signal-to-Noise Ratio",
        "Measures pixel-level distortion after embed",
        "Target: PSNR > 40 dB (imperceptible change)",
        "SSIM — Structural Similarity Index",
        "Compares luminance, contrast & structure",
        "Target: SSIM > 0.99",
    ]),
    ("🔐  Security",YELLOW,[
        "Histogram comparison — original vs stego",
        "Near-identical curves = low detectability",
        "Chi-square test — pixel pair randomness",
        "SPAM feature analysis resistance",
        "Goal: statistical indistinguishability",
    ]),
    ("⚡  Performance",ACCENT2,[
        "Payload capacity vs image distortion trade-off",
        "Accuracy of message extraction (100% target)",
        "Embedding speed on various image sizes",
        "Resistance to image processing (compress/resize)",
        "Benchmarked on standard test images",
    ]),
]

for i,(title,col,items) in enumerate(evals):
    lx=Inches(0.4+i*4.35)
    card(sl,lx,Inches(1.4),Inches(4.1),Inches(5.7))
    rect(sl,lx,Inches(1.4),Inches(4.1),Inches(0.5),col)
    txt(sl,title,lx+Inches(0.1),Inches(1.4),Inches(3.9),Inches(0.5),
        size=17,bold=True,col=BG)
    bullets(sl,items,lx+Inches(0.15),Inches(2.0),Inches(3.8),Inches(4.9),size=14)

# ════════════════════════════════════════════════════════════════════════
# SLIDE I – References
# ════════════════════════════════════════════════════════════════════════
sl=prs.slides.add_slide(blank); add_bg(sl)
header(sl,"References","Cited Literature & Resources")

refs=[
    ("[1] Pevný, T., Bas, P., Fridrich, J. (2010). Steganalysis by Subtractive Pixel Adjacency Matrix. IEEE TIFS.",
     "https://ws2.binghamton.edu/fridrich/Research/paper_6_dc.pdf"),
    ("[2] Fridrich, J., Kodovský, J. (2012). Rich Models for Steganalysis of Digital Images. IEEE TIFS.",
     "https://dde.binghamton.edu/kodovsky/pdf/TIFS2012-SRM.pdf"),
    ("[3] Yang, Z., Zhang, Y., Li, Y. (2023). IStego100K: A Large-Scale Image Steganalysis Dataset. IEEE Access.",
     "https://ieeexplore.ieee.org/document/6249077"),
    ("[4] Holub, V., Fridrich, J., Denemark, T. (2014). Universal Distortion Function for Steganography. IEEE TIFS.",
     "https://www.researchgate.net/publication/259639875"),
    ("[5] Apau, R., Badu, K., Amankwah, R. (2025). Image Steganography Techniques for Resisting Statistical Steganalysis. Scientific Reports.",
     "https://pmc.ncbi.nlm.nih.gov/articles/PMC11404826/"),
    ("[6] De La Croix, N.J. (2024). Comprehensive Survey on Image Steganalysis Using Deep Learning. Results in Engineering.",
     "https://www.sciencedirect.com/science/article/pii/S2590005624000195"),
]

for i,(ref,url) in enumerate(refs):
    ty=Inches(1.4+i*0.96)
    card(sl,Inches(0.4),ty,Inches(12.5),Inches(0.82))
    txt(sl,ref,Inches(0.6),ty+Inches(0.02),Inches(12.0),Inches(0.45),
        size=13,col=WHITE,bold=False)
    txt(sl,url,Inches(0.6),ty+Inches(0.44),Inches(12.0),Inches(0.35),
        size=11,col=ACCENT,italic=True)

# ════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════
out="StegAI_SDP_Presentation.pptx"
prs.save(out)
print(f"\n✅  Done! Saved → {out}")
print(f"   Total slides now: {len(prs.slides)}")
